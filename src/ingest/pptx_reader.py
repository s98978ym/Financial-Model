"""PPTX (PowerPoint) document text extraction with slide numbers preserved.

Uses python-pptx to iterate over slides and extract text from every shape
that carries a text frame (titles, body placeholders, free text boxes,
grouped shapes) as well as structured table data.
"""
import logging
from pathlib import Path
from typing import List

from .base import DocumentContent, PageContent

logger = logging.getLogger(__name__)


def _extract_table_from_shape(shape) -> List[List[str]]:
    """Return rows of cell text from a table shape."""
    rows: List[List[str]] = []
    try:
        table = shape.table
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            rows.append(cells)
    except Exception as exc:
        logger.warning("Failed to extract table from shape: %s", exc)
    return rows


def _collect_texts_from_shape(shape) -> List[str]:
    """Recursively collect text fragments from a shape (including groups)."""
    texts: List[str] = []

    # Grouped shapes — recurse into each child.
    try:
        if shape.shape_type is not None:
            from pptx.enum.shapes import MSO_SHAPE_TYPE  # type: ignore
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                for child_shape in shape.shapes:
                    texts.extend(_collect_texts_from_shape(child_shape))
                return texts
    except Exception:
        pass

    # Text frame (titles, body, free text boxes, etc.)
    if shape.has_text_frame:
        try:
            for paragraph in shape.text_frame.paragraphs:
                para_text = paragraph.text.strip()
                if para_text:
                    texts.append(para_text)
        except Exception as exc:
            logger.debug("Could not read text frame: %s", exc)

    return texts


def extract_pptx(file_path: str) -> DocumentContent:
    """Extract text and tables from a PPTX file, one page per slide.

    Parameters
    ----------
    file_path : str
        Path to the PPTX file.

    Returns
    -------
    DocumentContent
        Extracted document with per-slide text and tables.

    Raises
    ------
    FileNotFoundError
        If *file_path* does not exist.
    ImportError
        If python-pptx is not installed.
    RuntimeError
        If the file cannot be read.
    """
    path = Path(file_path)
    if not path.exists():
        raise FileNotFoundError(f"PPTX file not found: {file_path}")
    if not path.is_file():
        raise ValueError(f"Path is not a file: {file_path}")

    try:
        from pptx import Presentation  # type: ignore
    except ImportError:
        raise ImportError(
            "python-pptx is required for PPTX extraction. "
            "Install it with `pip install python-pptx`."
        )

    try:
        prs = Presentation(file_path)
    except Exception as exc:
        raise RuntimeError(
            f"Failed to open PPTX file: {file_path} — {exc}"
        ) from exc

    # ------------------------------------------------------------------
    # Metadata
    # ------------------------------------------------------------------
    metadata: dict = {}
    try:
        cp = prs.core_properties
        if cp.title:
            metadata["Title"] = cp.title
        if cp.author:
            metadata["Author"] = cp.author
        if cp.subject:
            metadata["Subject"] = cp.subject
        if cp.created:
            metadata["Created"] = str(cp.created)
        if cp.modified:
            metadata["Modified"] = str(cp.modified)
    except Exception:
        pass

    # ------------------------------------------------------------------
    # Iterate slides
    # ------------------------------------------------------------------
    pages: List[PageContent] = []
    total_slides = len(prs.slides)

    for slide_idx, slide in enumerate(prs.slides):
        slide_number = slide_idx + 1
        slide_texts: List[str] = []
        slide_tables: List[List[List[str]]] = []

        for shape in slide.shapes:
            # --- tables -----------------------------------------------
            if shape.has_table:
                table_rows = _extract_table_from_shape(shape)
                if table_rows:
                    slide_tables.append(table_rows)
                    # Also add a textual representation so full_text search
                    # finds table contents.
                    for row in table_rows:
                        slide_texts.append(" | ".join(row))
                continue  # table shapes usually don't carry extra text

            # --- text shapes ------------------------------------------
            texts = _collect_texts_from_shape(shape)
            slide_texts.extend(texts)

        # Build slide notes if present.
        if slide.has_notes_slide:
            try:
                notes_frame = slide.notes_slide.notes_text_frame
                notes_text = notes_frame.text.strip()
                if notes_text:
                    slide_texts.append(f"[Speaker Notes]\n{notes_text}")
            except Exception:
                pass

        pages.append(
            PageContent(
                page_number=slide_number,
                text="\n".join(slide_texts),
                tables=slide_tables,
                source_type="pptx",
            )
        )

    # Handle empty presentations.
    if not pages:
        pages.append(
            PageContent(
                page_number=1,
                text="",
                tables=[],
                source_type="pptx",
            )
        )
        total_slides = 1

    return DocumentContent(
        file_path=str(file_path),
        file_type="pptx",
        pages=pages,
        total_pages=total_slides,
        metadata=metadata,
    )
